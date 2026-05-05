import streamlit as st
import pandas as pd
import plotly.graph_objects as go
from pathlib import Path
from pptx import Presentation

# ─── CONFIGURACIÓN ─────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Reporte de Visitas Comerciales",
    page_icon="📊",
    layout="wide",
)

# ─── ESTILOS ───────────────────────────────────────────────────────────────────
st.markdown("""
<style>
    .kpi-container {
        display: flex;
        justify-content: space-between;
        gap: 20px;
        margin-top: 10px;
        margin-bottom: 25px;
    }
    .kpi-card {
        flex: 1;
        background-color: var(--secondary-background-color) !important;
        border-radius: 8px;
        padding: 24px 20px;
        border-left: 4px solid #1c64f2; 
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
        color: var(--text-color) !important;
        position: relative;
    }
    .kpi-card-icon {
        position: absolute;
        top: 20px;
        right: 20px;
        width: 36px;
        height: 36px;
        border-radius: 50%;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 18px;
    }
    .kpi-card:nth-child(1) { border-left-color: #3b82f6; }
    .kpi-card:nth-child(1) .kpi-card-icon { background-color: #dbeafe; color: #1d4ed8; }
    .kpi-card:nth-child(2) { border-left-color: #10b981; }
    .kpi-card:nth-child(2) .kpi-card-icon { background-color: #d1fae5; color: #047857; }
    .kpi-card:nth-child(3) { border-left-color: #f59e0b; }
    .kpi-card:nth-child(3) .kpi-card-icon { background-color: #fef3c7; color: #b45309; }
    .kpi-card:nth-child(4) { border-left-color: #8b5cf6; }
    .kpi-card:nth-child(4) .kpi-card-icon { background-color: #ede9fe; color: #6d28d9; }
    .kpi-value {
        font-size: 32px;
        font-weight: 800;
        margin-top: 10px;
        margin-bottom: 5px;
        color: #1e3a8a;
    }
    .kpi-label { font-size: 14px; opacity: 0.7; }

    .dashboard-panel {
        background-color: var(--secondary-background-color) !important;
        border-radius: 12px;
        padding: 24px;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
        color: var(--text-color) !important;
        border: 1px solid rgba(128,128,128,0.2);
        height: 100%;
    }
    .panel-title { font-size: 18px; font-weight: 700; color: #1e3a8a; margin-bottom: 4px; }
    .panel-subtitle { font-size: 13px; opacity: 0.6; margin-bottom: 20px; }

    .rutas-container { display: flex; flex-wrap: wrap; gap: 15px; }
    .ruta-card {
        flex: 1 1 calc(50% - 15px);
        min-width: 130px;
        border: 1px solid rgba(128,128,128,0.2);
        border-radius: 10px;
        padding: 15px;
    }
    .ruta-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 12px; }
    .ruta-name { font-size: 13px; opacity: 0.8; }
    .ruta-value { font-weight: 700; font-size: 15px; color: #1e3a8a; }
    .ruta-progress-container { display: flex; align-items: center; gap: 10px; }
    .ruta-progress-bar {
        flex-grow: 1;
        background-color: rgba(128,128,128,0.2);
        height: 6px;
        border-radius: 3px;
        overflow: hidden;
    }
    .ruta-progress-fill { height: 100%; border-radius: 3px; }
    .ruta-pct { font-size: 12px; opacity: 0.6; min-width: 40px; text-align: right; }

    .conv-list { display: flex; flex-direction: column; gap: 16px; }
    .conv-item { display: flex; justify-content: space-between; align-items: center; }
    .conv-item-left { display: flex; align-items: center; gap: 10px; }
    .conv-dot { width: 8px; height: 8px; border-radius: 50%; }
    .conv-name { font-size: 14px; }
    .conv-value { font-weight: 700; font-size: 14px; }

    .styled-table { width: 100%; border-collapse: collapse; font-size: 13px; }
    .styled-table thead tr {
        border-bottom: 2px solid rgba(128,128,128,0.2);
        color: #1e3a8a;
        text-align: left;
    }
    .styled-table th, .styled-table td { padding: 10px 8px; }
    .styled-table tbody tr { border-bottom: 1px solid rgba(128,128,128,0.1); }
    .styled-table tbody tr:last-of-type { border-bottom: none; }

    .section-header {
        border-radius: 8px;
        padding: 14px 20px;
        margin-bottom: 16px;
        display: flex;
        justify-content: space-between;
        align-items: center;
        color: white;
    }
    .section-header-title {
        font-size: 19px;
        font-weight: 800;
        display: flex;
        align-items: center;
        gap: 10px;
    }
    .section-header-sub { font-size: 13px; opacity: 0.85; margin-top: 3px; }
    .section-header-right { font-size: 13px; font-weight: 600; text-align: right; }

    .header-title-container-main { text-align: center; margin-top: 10px; margin-bottom: 15px; }
    .header-title { color: #1e3a8a; font-size: 38px; font-weight: 800; margin-bottom: 5px; }
</style>
""", unsafe_allow_html=True)

_theme   = st.get_option("theme.base") or "dark"
_is_dark = (_theme == "dark")

if _is_dark:
    st.markdown("""
    <style>
        .header-title, .panel-title, .kpi-value, .ruta-value { color: #60a5fa !important; }
        .styled-table thead tr { color: #60a5fa !important; }
    </style>
    """, unsafe_allow_html=True)

# ─── CONSTANTES ────────────────────────────────────────────────────────────────
EXCEL_FILE = "visitas_ventas.xlsx"

COL_FECHA    = "Date"
COL_TIPO     = "Tipo"
COL_TIPO_CLI = "Giro"
COL_CLIENTE  = "Cliente o Prospecto"
COL_DISTRITO = "Distrito"
COL_MOTIVO   = "Task"
COL_RESULTADO= "Obs"
COL_ZONA     = "Zona"
COL_REGION   = "Región"
COL_TIPO_VIS = "Tipo Visita"
COL_VENDEDOR = "Vendedor"

ETAPAS_EMBUDO = [
    "PROSPECCIÓN", "CALIFICACIÓN DE LEADS", "VISITA",
    "PROPUESTA", "NEGOCIACIÓN", "CIERRE", "NO CIERRE",
]
PALETA_RUTAS = ["#3b82f6","#10b981","#f59e0b","#8b5cf6","#ef4444","#0ea5e9","#14b8a6","#f43f5e"]

MESES_ES = {
    1:"Enero",2:"Febrero",3:"Marzo",4:"Abril",5:"Mayo",6:"Junio",
    7:"Julio",8:"Agosto",9:"Septiembre",10:"Octubre",11:"Noviembre",12:"Diciembre"
}

# ─── UTILIDADES DE FECHA ───────────────────────────────────────────────────────
def mes_a_label(ym_str):
    """'2025-06' → 'Junio 2025'"""
    try:
        anio, mes = ym_str.split("-")
        return f"{MESES_ES[int(mes)]} {anio}"
    except Exception:
        return ym_str

def sem_a_label(sem_str):
    """'2026-S17' → '2026-04-20 al 2026-04-26'"""
    try:
        anio, s = sem_str.split("-S")
        lunes = pd.Timestamp.fromisocalendar(int(anio), int(s), 1)
        domingo = lunes + pd.Timedelta(days=6)
        return f"{lunes.date()} al {domingo.date()}"
    except Exception:
        return sem_str

def label_a_semkey(label):
    """Reverse lookup: '2026-04-20 al 2026-04-26' → '2026-S17'"""
    # stored in map built at filter time
    return label

# ─── CARGA DE DATOS ────────────────────────────────────────────────────────────
@st.cache_data(ttl=300)
def cargar_datos_completos(ruta):
    df_log   = pd.read_excel(ruta, sheet_name="Log",   engine="openpyxl")
    df_users = pd.read_excel(ruta, sheet_name="Users", engine="openpyxl")
    df_zona  = pd.read_excel(ruta, sheet_name="Zona",  engine="openpyxl")

    for df in [df_log, df_users, df_zona]:
        df.columns = df.columns.str.strip()

    if "User" in df_log.columns and "Email" in df_users.columns:
        df_log = df_log.merge(df_users[["Email","Name"]], left_on="User", right_on="Email", how="left")
        df_log[COL_VENDEDOR] = df_log["Name"].fillna("Desconocido")
    else:
        df_log[COL_VENDEDOR] = "Desconocido"

    if "Zona" in df_log.columns and "Zona" in df_zona.columns:
        df_log = df_log.merge(df_zona[["Zona","Tipo Zona"]], on="Zona", how="left")
        df_log[COL_REGION] = df_log["Tipo Zona"].fillna("Desconocido")
    else:
        df_log[COL_REGION] = "Desconocido"

    df_log[COL_FECHA] = pd.to_datetime(df_log[COL_FECHA], dayfirst=True, errors="coerce")
    df_log = df_log.dropna(subset=[COL_FECHA])

    iso = df_log[COL_FECHA].dt.isocalendar()
    df_log["_sem_key"] = iso["year"].astype(str) + "-S" + iso["week"].astype(str).str.zfill(2)
    df_log["_mes_key"] = df_log[COL_FECHA].dt.strftime("%Y-%m")

    for col in [COL_VENDEDOR,COL_TIPO,COL_TIPO_CLI,COL_CLIENTE,COL_DISTRITO,COL_MOTIVO,COL_TIPO_VIS,COL_REGION,COL_ZONA]:
        if col in df_log.columns:
            df_log[col] = df_log[col].astype(str).str.strip()

    try:
        df_est_sem = pd.read_excel(ruta, sheet_name="Estado_Semana", engine="openpyxl")
    except Exception:
        df_est_sem = pd.DataFrame(columns=["Estado","Cantidad"])

    try:
        df_est_mes = pd.read_excel(ruta, sheet_name="Estado_Mes", engine="openpyxl")
    except Exception:
        df_est_mes = pd.DataFrame(columns=["Estado","Cantidad"])

    # ── Limpieza: eliminar filas con campos obligatorios vacíos ──────────────
    CAMPOS_OBLIGATORIOS = [
        "Id", "Date", "Time", "Zona", "Cliente o Prospecto",
        "Tipo", "Task", "User", "Giro", "Departamento",
        "Provincia", "Distrito", "Tipo Visita", "Id Cliente Prospecto",
    ]
    cols_a_validar = [c for c in CAMPOS_OBLIGATORIOS if c in df_log.columns]
    filas_antes = len(df_log)

    for c in cols_a_validar:
        try:
            if df_log[c].dtype == object:
                # Convertir a str, limpiar espacios, marcar vacíos como NaN
                df_log[c] = (df_log[c].astype(str).str.strip()
                             .replace({"nan": None, "None": None,
                                       "NaT": None, "": None}))
        except Exception:
            pass  # columnas numéricas/fecha se dejan para dropna

    df_log = df_log.dropna(subset=cols_a_validar)
    df_log._filas_descartadas = filas_antes - len(df_log)

    return df_log, df_est_sem, df_est_mes

# ── Selector de archivo ───────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### ⚙️ Configuración")
    archivo = st.file_uploader("Cargar Excel de visitas", type=["xlsx","xls","xlsm"])

    if archivo is None:
        st.warning("⚠️ Sube un archivo Excel para comenzar.")
        st.stop()

    try:
        df_raw, df_est_sem, df_est_mes = cargar_datos_completos(archivo)
    except Exception as e:
        st.error(f"Error al leer el archivo: {e}")
        st.stop()

    if df_raw.empty:
        st.error("El archivo está vacío o no tiene filas válidas.")
        st.stop()

    df_raw = df_raw.sort_values(COL_FECHA)

    st.divider()
    st.markdown("### 🔽 Filtros")

    vendedores_list = ["Todos"] + sorted(df_raw[COL_VENDEDOR].dropna().unique().tolist())
    sel_vendedor = st.selectbox("Vendedor", vendedores_list)

    st.markdown("#### Periodo de Filtro")
    import datetime as _dt
    _hoy = _dt.date.today()
    _anio_hoy = str(_hoy.year)
    _mes_hoy  = _hoy.month
    _iso      = _hoy.isocalendar()
    _sem_hoy  = f"{_iso[0]}-S{str(_iso[1]).zfill(2)}"
    # Inicializar solo en primera carga
    if "_filtro_inicializado" not in st.session_state:
        st.session_state["modo_fecha_radio"] = "Semana"
        st.session_state["anio_ini"] = _anio_hoy
        st.session_state["mes_ini"]  = _mes_hoy
        st.session_state["anio_fin"] = _anio_hoy
        st.session_state["mes_fin"]  = _mes_hoy
        st.session_state["sem_ini"]  = _sem_hoy
        st.session_state["sem_fin"]  = _sem_hoy
        st.session_state["_filtro_inicializado"] = True
    modo_fecha = st.radio("Agrupar por:", ["Mes", "Semana"], horizontal=True, key="modo_fecha_radio")

    # ── Metadatos de fechas disponibles en el archivo ─────────────────────────
    mes_keys_disp = sorted(df_raw["_mes_key"].dropna().unique().tolist())
    anios_disp    = sorted({k[:4] for k in mes_keys_disp})
    meses_por_anio = {}
    for k in mes_keys_disp:
        a, m = k[:4], int(k[5:7])
        meses_por_anio.setdefault(a, []).append(m)

    sem_keys_disp = sorted(df_raw["_sem_key"].dropna().unique().tolist())
    sem_por_mes_key = {}
    for sk in sem_keys_disp:
        anio_s, s_nro = sk.split("-S")
        lunes = pd.Timestamp.fromisocalendar(int(anio_s), int(s_nro), 1)
        mk    = lunes.strftime("%Y-%m")
        sem_por_mes_key.setdefault(mk, []).append(sk)

    def sem_display(sk):
        """Texto para los selectboxes: solo rango de fechas con separador visual."""
        anio_s, s_nro_s = sk.split("-S")
        lunes   = pd.Timestamp.fromisocalendar(int(anio_s), int(s_nro_s), 1)
        domingo = lunes + pd.Timedelta(days=6)
        return f"{lunes.date()}  ──  {domingo.date()}"

    def sem_ini_label_narrativo(sk):
        """Fecha de inicio de semana para texto narrativo del dashboard."""
        anio_s, s_nro_s = sk.split("-S")
        lunes = pd.Timestamp.fromisocalendar(int(anio_s), int(s_nro_s), 1)
        return str(lunes.date())

    def sem_fin_label_narrativo(sk):
        """Fecha de fin de semana para texto narrativo del dashboard."""
        anio_s, s_nro_s = sk.split("-S")
        lunes   = pd.Timestamp.fromisocalendar(int(anio_s), int(s_nro_s), 1)
        domingo = lunes + pd.Timedelta(days=6)
        return str(domingo.date())

    # ── INICIO ────────────────────────────────────────────────────────────────
    st.markdown("**📅 Inicio**")
    col_ai, col_mi = st.columns(2)
    with col_ai:
        anio_ini = st.selectbox("Año", anios_disp, index=0, key="anio_ini")
    with col_mi:
        meses_ini_disp = meses_por_anio.get(anio_ini, [])
        # Si el año cambió, resetear mes_ini al primero disponible
        if st.session_state.get("_prev_anio_ini") != anio_ini:
            st.session_state["mes_ini"] = meses_ini_disp[0] if meses_ini_disp else None
        mes_ini = st.selectbox(
            "Mes", meses_ini_disp,
            format_func=lambda m: MESES_ES[m],
            key="mes_ini"
        )
    ini_mes_key = f"{anio_ini}-{str(mes_ini).zfill(2)}"

    if modo_fecha == "Semana":
        sems_ini = sem_por_mes_key.get(ini_mes_key, [])
        if sems_ini:
            # Si el mes_ini cambió, resetear semana inicio
            if st.session_state.get("_prev_mes_ini") != mes_ini or st.session_state.get("_prev_anio_ini") != anio_ini:
                st.session_state["sem_ini"] = sems_ini[0]
            sem_ini_sk = st.selectbox(
                "Semana inicio", sems_ini,
                format_func=sem_display,
                key="sem_ini"
            )
        else:
            st.caption("Sin semanas registradas en ese mes.")
            sem_ini_sk = sem_keys_disp[0]
    else:
        sem_ini_sk = sem_keys_disp[0]  # valor dummy cuando no aplica

    # ── FIN: auto-setear al inicio sólo cuando el inicio cambia ──────────────
    # Detectar si el inicio acaba de cambiar en esta ejecución
    _ini_cambio_anio = st.session_state.get("_prev_anio_ini") != anio_ini
    _ini_cambio_mes  = st.session_state.get("_prev_mes_ini")  != mes_ini
    _ini_cambio_sem  = st.session_state.get("_prev_sem_ini")  != sem_ini_sk

    st.markdown("**📅 Fin**")
    col_af, col_mf = st.columns(2)
    with col_af:
        if _ini_cambio_anio:
            st.session_state["anio_fin"] = anio_ini
        anio_fin = st.selectbox("Año ", anios_disp, key="anio_fin")
    with col_mf:
        meses_fin_disp = meses_por_anio.get(anio_fin, [])
        # Resetear mes_fin si anio_fin cambió por seguir al inicio, o si mes_ini cambió
        if _ini_cambio_anio or _ini_cambio_mes:
            nuevo_mes_fin = mes_ini if mes_ini in meses_fin_disp else (meses_fin_disp[0] if meses_fin_disp else None)
            st.session_state["mes_fin"] = nuevo_mes_fin
        # Si el año fin fue cambiado por el usuario y el mes guardado ya no es válido, resetear
        if st.session_state.get("mes_fin") not in meses_fin_disp and meses_fin_disp:
            st.session_state["mes_fin"] = meses_fin_disp[0]
        mes_fin = st.selectbox(
            "Mes ", meses_fin_disp,
            format_func=lambda m: MESES_ES[m],
            key="mes_fin"
        )
    fin_mes_key = f"{anio_fin}-{str(mes_fin).zfill(2)}"

    if modo_fecha == "Semana":
        sems_fin = sem_por_mes_key.get(fin_mes_key, [])
        if sems_fin:
            if _ini_cambio_anio or _ini_cambio_mes or _ini_cambio_sem:
                st.session_state["sem_fin"] = sem_ini_sk if sem_ini_sk in sems_fin else sems_fin[0]
            if st.session_state.get("sem_fin") not in sems_fin:
                st.session_state["sem_fin"] = sems_fin[0]
            sem_fin_sk = st.selectbox(
                "Semana fin", sems_fin,
                format_func=sem_display,
                key="sem_fin"
            )
        else:
            st.caption("Sin semanas registradas en ese mes.")
            sem_fin_sk = sem_keys_disp[-1]
    else:
        sem_fin_sk = sem_keys_disp[-1]  # valor dummy

    # Guardar estado del inicio para detectar cambios en el próximo rerun
    st.session_state["_prev_anio_ini"] = anio_ini
    st.session_state["_prev_mes_ini"]  = mes_ini
    st.session_state["_prev_sem_ini"]  = sem_ini_sk

    # ── Claves y labels finales ────────────────────────────────────────────────
    if modo_fecha == "Mes":
        opciones_keys = mes_keys_disp
        ini_key = ini_mes_key
        fin_key = fin_mes_key
        if ini_key > fin_key:
            ini_key, fin_key = fin_key, ini_key
        sel_ini_key, sel_fin_key = ini_key, fin_key
        # Labels para texto narrativo del dashboard
        sel_ini_label = mes_a_label(ini_key)   # "Junio 2025"
        sel_fin_label = mes_a_label(fin_key)
        COL_FILTRO = "_mes_key"
    else:
        opciones_keys = sem_keys_disp
        ini_key = sem_ini_sk
        fin_key = sem_fin_sk
        if ini_key > fin_key:
            ini_key, fin_key = fin_key, ini_key
        sel_ini_key, sel_fin_key = ini_key, fin_key
        # Para semanas: etiqueta narrativa = "<lunes ini> y <domingo fin>"
        sel_ini_label = sem_ini_label_narrativo(ini_key)   # "2026-04-13" (lunes inicio)
        sel_fin_label = sem_fin_label_narrativo(fin_key)   # "2026-04-26" (domingo fin)
        COL_FILTRO = "_sem_key"

    st.divider()
    if st.button("↻ Limpiar caché"):
        st.cache_data.clear()
        st.rerun()




# ── Aplicar filtros ────────────────────────────────────────────────────────────
dff = df_raw.copy()
if sel_vendedor != "Todos":
    dff = dff[dff[COL_VENDEDOR] == sel_vendedor]

dff = dff[(dff[COL_FILTRO] >= sel_ini_key) & (dff[COL_FILTRO] <= sel_fin_key)]

# Excluir zona OFICINA de todos los analisis
if COL_ZONA in dff.columns:
    dff = dff[dff[COL_ZONA].str.upper() != 'OFICINA']

# num_periodos = cantidad de periodos en el rango
idx_ini = opciones_keys.index(sel_ini_key)
idx_fin = opciones_keys.index(sel_fin_key)
num_periodos = abs(idx_fin - idx_ini) + 1

# ── Etiqueta legible del rango ─────────────────────────────────────────────────
if modo_fecha == "Mes":
    if sel_ini_key == sel_fin_key:
        rango_label = f"en el mes de {sel_ini_label}"
    else:
        rango_label = f"entre {sel_ini_label} y {sel_fin_label}"
else:  # Semana
    if sel_ini_key == sel_fin_key:
        # Una sola semana: lunes al domingo
        rango_label = f"entre {sel_ini_label} y {sel_fin_label}"
    else:
        rango_label = f"entre {sel_ini_label} y {sel_fin_label}"

df_estado_usar = df_est_mes if modo_fecha == "Mes" else df_est_sem

with st.sidebar:
    st.caption(f"**{len(dff):,}** registros · **{num_periodos}** {'meses' if modo_fecha=='Mes' else 'semana(s)'}")

# ─── HEADER PRINCIPAL ─────────────────────────────────────────────────────────
st.markdown("""
<div class="header-title-container-main">
    <div class="header-title">🗺️ Reporte de Visitas Comerciales</div>
</div>
""", unsafe_allow_html=True)
st.markdown(
    f"<div style='text-align:center;margin-bottom:2rem;opacity:0.7;font-size:15px;'>"
    f"Análisis de rutas, conversión y cobertura comercial {rango_label}</div>",
    unsafe_allow_html=True
)

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS DE CÁLCULO
# ═══════════════════════════════════════════════════════════════════════════════
def calc_kpis(df_filtro):
    if COL_TIPO not in df_filtro.columns:
        return 0, 0, 0, 0.0
    df_vis  = df_filtro[
        df_filtro[COL_TIPO].str.upper().isin(["PROSPECCIÓN","PROSPECCION","MANTENIMIENTO"]) &
        df_filtro.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
    ]
    tot_vis = len(df_vis)
    df_pros = df_filtro[df_filtro[COL_TIPO].str.upper().isin(["PROSPECCIÓN","PROSPECCION"])]
    tot_pros = df_pros[COL_CLIENTE].nunique() if not df_pros.empty else 0
    n_cierres = 0
    if not df_pros.empty and COL_MOTIVO in df_pros.columns:
        orden = {e: i for i, e in enumerate(ETAPAS_EMBUDO)}
        valid = df_pros[df_pros[COL_MOTIVO].str.upper().isin([e.upper() for e in ETAPAS_EMBUDO])].copy()
        if not valid.empty:
            valid["_ord"] = valid[COL_MOTIVO].str.upper().map(orden)
            ultima = valid.sort_values("_ord").groupby(COL_CLIENTE).last()[[COL_MOTIVO]].reset_index()
            n_cierres = ultima[ultima[COL_MOTIVO].str.upper() == "CIERRE"].shape[0]
    t_conv = round(n_cierres / tot_pros * 100, 2) if tot_pros > 0 else 0.0
    return tot_vis, tot_pros, n_cierres, t_conv


def calc_visitas_planificadas(df_region, meta_df, num_p):
    # Visitas planificadas = umbral_activo x zonas_activas x num_periodos
    if COL_ZONA not in df_region.columns:
        return None
    df_fis = df_region[
        df_region[COL_TIPO].str.upper().isin(['PROSPECCION', 'PROSPECCIÓN', 'MANTENIMIENTO']) &
        df_region.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(['FISICA', 'FÍSICA'])
    ]
    zonas_activas = df_fis[COL_ZONA].nunique()
    if zonas_activas == 0:
        return None
    if meta_df.empty or 'Cantidad' not in meta_df.columns or 'Estado' not in meta_df.columns:
        return None
    meta_cp = meta_df.copy()
    meta_cp['Cantidad'] = pd.to_numeric(meta_cp['Cantidad'], errors='coerce').fillna(0)
    fila_activo = meta_cp[meta_cp['Estado'].astype(str).str.upper().str.contains('ACTIVO', na=False)]
    if fila_activo.empty:
        return None
    umbral = int(fila_activo['Cantidad'].iloc[0])
    return umbral * zonas_activas * num_p


def obtener_estado(visitas, meta_df, num_p):
    if meta_df.empty or "Cantidad" not in meta_df.columns or "Estado" not in meta_df.columns:
        return "Sin Datos"
    meta_df = meta_df.copy()
    meta_df["Cantidad"] = pd.to_numeric(meta_df["Cantidad"], errors="coerce").fillna(0)
    mdf = meta_df.sort_values("Cantidad", ascending=False)
    for _, row in mdf.iterrows():
        if visitas >= row["Cantidad"] * num_p:
            return str(row["Estado"]).strip()
    return str(mdf.iloc[-1]["Estado"]).strip()


def color_estado(est):
    eu = str(est).upper()
    if "ACTIVO" in eu:   return "#10b981"
    if "REGULAR" in eu:  return "#f59e0b"
    if "BAJO" in eu:     return "#ef4444"
    return "var(--text-color)"


def build_rutas_html(grupos, total):
    html = ""
    for i, row in grupos.iterrows():
        nombre = row[COL_ZONA]
        vis    = row["Visitas"]
        pct    = round(vis / total * 100, 2) if total > 0 else 0
        color  = PALETA_RUTAS[i % len(PALETA_RUTAS)]
        html += f"""<div class="ruta-card">
<div class="ruta-header"><span class="ruta-name">{nombre}</span><span class="ruta-value">{vis}</span></div>
<div class="ruta-progress-container">
<div class="ruta-progress-bar"><div class="ruta-progress-fill" style="width:{pct}%;background:{color};"></div></div>
<span class="ruta-pct">{pct}%</span></div></div>"""
    return html


def build_tabla_estado_html(grupos, total):
    html = """<table class="styled-table">
<thead><tr><th>Zona</th><th>Visitas</th><th>%</th><th>Estado</th></tr></thead><tbody>"""
    for i, row in grupos.iterrows():
        zona = row[COL_ZONA]
        vis  = row["Visitas"]
        pct  = round(vis / total * 100, 2) if total > 0 else 0
        est  = obtener_estado(vis, df_estado_usar, num_periodos)
        col  = color_estado(est)
        html += f"""<tr>
<td style="font-weight:700;">{zona}</td>
<td>{vis}</td>
<td>{pct}%</td>
<td style="color:{col};font-weight:700;">{est}</td></tr>"""
    html += "</tbody></table>"
    return html


def seccion_header(icono, titulo, subtitulo, info_derecha, gradient="linear-gradient(135deg,#1d4ed8,#3b82f6)"):
    st.markdown(f"""
<div class="section-header" style="background:{gradient};">
<div>
<div class="section-header-title">{icono} {titulo}</div>
<div class="section-header-sub">{subtitulo}</div>
</div>
<div class="section-header-right">{info_derecha}</div>
</div>""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# RENDER PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════
def render_region_dashboard(df_region, region_nombre, is_todos=False):
    # ── CABECERA + KPIs – siempre primero ─────────────────────────────────────
    seccion_header(
        "📊", "Resumen Ejecutivo",
        f"Análisis de visitas comerciales · {region_nombre} · {rango_label}",
        f"Actualizado: {sel_fin_label}",
        gradient="linear-gradient(135deg,#1d4ed8 0%,#3b82f6 100%)"
    )

    tot_vis, tot_pros, n_cierres, t_conv = calc_kpis(df_region)
    vis_plan = calc_visitas_planificadas(df_region, df_estado_usar, num_periodos)
    if vis_plan and vis_plan > 0:
        tasa_vis = round(tot_vis / vis_plan * 100, 1)
        vis_kpi_val = f'{tot_vis:,} / {vis_plan:,}'
        vis_kpi_sub = f'<div style="font-size:13px;color:#10b981;font-weight:700;margin-top:2px;">{tasa_vis}% del planificado</div>'
    else:
        vis_kpi_val = str(tot_vis)
        vis_kpi_sub = ''
    st.markdown(f"""<div class="kpi-container">
<div class="kpi-card"><div class="kpi-label">Visitas Totales</div><div class="kpi-value">{vis_kpi_val}</div>{vis_kpi_sub}<div class="kpi-card-icon">👥</div></div>
<div class="kpi-card"><div class="kpi-label">Prospectos Únicos</div><div class="kpi-value">{tot_pros}</div><div class="kpi-card-icon">👤</div></div>
<div class="kpi-card"><div class="kpi-label">Cierres</div><div class="kpi-value">{n_cierres}</div><div class="kpi-card-icon">✔️</div></div>
<div class="kpi-card"><div class="kpi-label">Conversión</div><div class="kpi-value">{t_conv}%</div><div class="kpi-card-icon">%</div></div>
</div>""", unsafe_allow_html=True)

    if is_todos:
        return

    if COL_ZONA not in df_region.columns:
        st.warning("No se cuenta con la columna 'Zona'.")
        return

    # ── DATOS SECCIÓN 1 (Prospección + Mantenimiento, Físicas) ────────────────
    df_fis = df_region[
        df_region[COL_TIPO].str.upper().isin(["PROSPECCIÓN","PROSPECCION","MANTENIMIENTO"]) &
        df_region.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
    ]
    grupos_s1 = df_fis.groupby(COL_ZONA).size().reset_index(name="Visitas").sort_values("Visitas", ascending=False).reset_index(drop=True)
    total_s1  = grupos_s1["Visitas"].sum()

    # ── Conversiones por zona (Prospección) ───────────────────────────────────
    zonas = df_region[COL_ZONA].dropna().unique()
    convs = []
    for i, z in enumerate(zonas):
        _, p, _, ct = calc_kpis(df_region[df_region[COL_ZONA] == z])
        if p > 0:
            convs.append({"zona": z, "tasa": ct, "color": PALETA_RUTAS[i % len(PALETA_RUTAS)]})
    convs.sort(key=lambda x: x["tasa"], reverse=True)

    html_convs = "".join(f"""<div class="conv-item">
<div class="conv-item-left"><div class="conv-dot" style="background:{c['color']};"></div><span class="conv-name">{c['zona']}</span></div>
<span class="conv-value">{c['tasa']}%</span></div>""" for c in convs) or \
        "<span style='opacity:0.6;font-size:13px;'>Sin datos de prospección suficientes.</span>"

    html_rutas_s1 = build_rutas_html(grupos_s1, total_s1)
    html_tabla_s1 = build_tabla_estado_html(grupos_s1, total_s1)

    # ── PANEL 2 COLUMNAS ──────────────────────────────────────────────────────
    col1, col2 = st.columns([1, 1])

    with col1:
        st.markdown("""<div class="dashboard-panel">
<div class="panel-title">Distribución General</div>
<div class="panel-subtitle">Visitas físicas (Prospección + Mantenimiento)</div>""", unsafe_allow_html=True)
        if html_rutas_s1:
            st.markdown(f'<div class="rutas-container" style="margin-top:14px;">{html_rutas_s1}</div>', unsafe_allow_html=True)
        else:
            st.markdown("<p style='opacity:0.6;font-size:13px;'>Sin registros.</p>", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with col2:
        st.markdown(f"""<div class="dashboard-panel">
<div class="panel-title">Conversiones</div>
<div class="panel-subtitle">% de cierres / prospectos por Zona</div>
<div class="conv-list">{html_convs}</div>
</div>""", unsafe_allow_html=True)

    st.markdown("<div style='margin-top:30px;'></div>", unsafe_allow_html=True)

    # ── SECCIÓN 2: MANTENIMIENTO ───────────────────────────────────────────────
    df_mant = df_region[
        df_region[COL_TIPO].str.upper() == "MANTENIMIENTO"
    ]
    df_mant_fis = df_mant[
        df_mant.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
    ]
    grupos_mant = df_mant_fis.groupby(COL_ZONA).size().reset_index(name="Visitas").sort_values("Visitas", ascending=False).reset_index(drop=True)
    total_mant  = grupos_mant["Visitas"].sum()

    seccion_header(
        "🔧", "Visitas por Zona - Mantenimiento",
        f"Distribución de visitas de mantenimiento físico por Zona · {region_nombre}",
        f"Total: {total_mant} visitas físicas",
        gradient="linear-gradient(135deg,#065f46 0%,#10b981 100%)"
    )

    if not grupos_mant.empty:
        html_tabla_mant = build_tabla_estado_html(grupos_mant, total_mant)
        col_graf, col_det = st.columns([3, 2])

        with col_graf:
            st.markdown("""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title">Distribución de Visitas por Zona</div>
<div class="panel-subtitle">Solo Mantenimiento · Tipo Visita = Física</div>
</div>""", unsafe_allow_html=True)
            fig = go.Figure(go.Bar(
                x=grupos_mant[COL_ZONA],
                y=grupos_mant["Visitas"],
                marker_color=[PALETA_RUTAS[i % len(PALETA_RUTAS)] for i in range(len(grupos_mant))],
                text=grupos_mant["Visitas"],
                textposition="outside",
                cliponaxis=False,
            ))
            fig.update_layout(
                paper_bgcolor="rgba(0,0,0,0)",
                plot_bgcolor="rgba(0,0,0,0)",
                margin=dict(t=10, b=10, l=10, r=10),
                xaxis=dict(showgrid=False),
                yaxis=dict(showgrid=True, gridcolor="rgba(128,128,128,0.2)"),
                height=320,
            )
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
            html_rutas_mant = build_rutas_html(grupos_mant, total_mant)
            st.markdown(f'<div class="rutas-container" style="margin-top:10px;">{html_rutas_mant}</div>', unsafe_allow_html=True)

        with col_det:
            st.markdown(f"""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title" style="margin-bottom:6px;">Detalle por Zona</div>
<div class="panel-subtitle" style="margin-bottom:10px;">Visitas, participación y estado · Prospección + Mantenimiento</div>
{html_tabla_s1}
</div>""", unsafe_allow_html=True)
    else:
        st.info("Sin registros de Mantenimiento físico para este periodo.")


# ═══════════════════════════════════════════════════════════════════════════════
# SECCIÓN: VISITAS POR ZONA - PROSPECCIÓN
# ═══════════════════════════════════════════════════════════════════════════════
def render_conversion_section(df_region, region_nombre):
    """Sección Visitas por Zona - Prospección (solo Tipo = PROSPECCIÓN, Tipo Visita = FÍSICA)."""

    # Filtrar: PROSPECCIÓN + FÍSICA
    df_pros_fis = df_region[
        df_region[COL_TIPO].str.upper().isin(["PROSPECCIÓN", "PROSPECCION"]) &
        df_region.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA", "FISICA"])
    ]
    if df_pros_fis.empty:
        return

    grupos_pros = (df_pros_fis.groupby(COL_ZONA).size()
                   .reset_index(name="Visitas")
                   .sort_values("Visitas", ascending=False)
                   .reset_index(drop=True))
    total_pros = grupos_pros["Visitas"].sum()

    st.markdown('<div style="margin-top:30px;"></div>', unsafe_allow_html=True)
    seccion_header(
        "🔍", "Visitas por Zona - Prospección",
        f"Distribución de visitas de prospección física por Zona · {region_nombre}",
        f"Total: {total_pros} visitas físicas",
        gradient="linear-gradient(135deg,#4c1d95 0%,#7c3aed 100%)"
    )

    html_tabla_pros = build_tabla_estado_html(grupos_pros, total_pros)
    col_graf, col_det = st.columns([3, 2])

    with col_graf:
        st.markdown("""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title">Distribución de Visitas por Zona</div>
<div class="panel-subtitle">Solo Prospección · Tipo Visita = Física</div>
</div>""", unsafe_allow_html=True)

        max_vis = int(grupos_pros["Visitas"].max()) if not grupos_pros.empty else 1
        fig = go.Figure(go.Bar(
            x=grupos_pros["Visitas"],
            y=grupos_pros[COL_ZONA],
            orientation="h",
            marker_color=[PALETA_RUTAS[i % len(PALETA_RUTAS)] for i in range(len(grupos_pros))],
            text=grupos_pros["Visitas"],
            textposition="outside",
            cliponaxis=False,
        ))
        fig.update_layout(
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            margin=dict(t=10, b=10, l=10, r=40),
            xaxis=dict(showgrid=True, gridcolor="rgba(128,128,128,0.15)",
                       range=[0, max_vis * 1.2]),
            yaxis=dict(showgrid=False, autorange="reversed"),
            height=max(220, len(grupos_pros) * 44),
        )
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

        html_rutas_pros = build_rutas_html(grupos_pros, total_pros)
        st.markdown(f'<div class="rutas-container" style="margin-top:10px;">{html_rutas_pros}</div>',
                    unsafe_allow_html=True)

    # ── Panel Alertas & Observaciones (conversión) ──────────────────────────────
    with col_det:
        # Calcular tasas de conversión por zona para las alertas
        df_pros_all = df_region[df_region[COL_TIPO].str.upper().isin(["PROSPECCIÓN", "PROSPECCION"])]
        orden_etapa = {e: i for i, e in enumerate(ETAPAS_EMBUDO)}
        zona_stats = []
        for i, z in enumerate(grupos_pros[COL_ZONA].tolist()):
            df_z = df_pros_all[df_pros_all[COL_ZONA] == z]
            n_pros = df_z[COL_CLIENTE].nunique()
            if n_pros == 0:
                continue
            n_cierres = 0
            if COL_MOTIVO in df_z.columns:
                valid = df_z[df_z[COL_MOTIVO].str.upper().isin([e.upper() for e in ETAPAS_EMBUDO])].copy()
                if not valid.empty:
                    valid["_ord"] = valid[COL_MOTIVO].str.upper().map(orden_etapa)
                    ultima = valid.sort_values("_ord").groupby(COL_CLIENTE).last()[[COL_MOTIVO]].reset_index()
                    n_cierres = ultima[ultima[COL_MOTIVO].str.upper() == "CIERRE"].shape[0]
            tasa = round(n_cierres / n_pros * 100, 1)
            zona_stats.append({"zona": z, "prospectos": n_pros, "cierres": n_cierres, "tasa": tasa})

        total_pros_gbl    = df_pros_all[COL_CLIENTE].nunique()
        total_cierres_gbl = sum(z["cierres"] for z in zona_stats)
        tasa_global = round(total_cierres_gbl / total_pros_gbl * 100, 2) if total_pros_gbl > 0 else 0.0

        sin_cierre = [z for z in zona_stats if z["tasa"] == 0 and z["prospectos"] > 0]
        top_zona   = max(zona_stats, key=lambda z: z["tasa"]) if zona_stats else None

        alertas_html = ""
        if sin_cierre:
            for z in sin_cierre:
                alertas_html += f"""
<div style="display:flex;align-items:flex-start;gap:10px;padding:10px 12px;
            background:rgba(239,68,68,0.08);border-left:3px solid #ef4444;
            border-radius:6px;margin-bottom:8px;">
<span style="font-size:16px;">&#9888;</span>
<div style="font-size:12px;">
<strong style="color:#ef4444;">{z['zona']}:</strong><br>
{z['prospectos']} prospectos, 0 cierres (0%)
</div></div>"""
        if top_zona and top_zona["tasa"] > 0:
            alertas_html += f"""
<div style="display:flex;align-items:flex-start;gap:10px;padding:10px 12px;
            background:rgba(16,185,129,0.08);border-left:3px solid #10b981;
            border-radius:6px;margin-bottom:8px;">
<span style="font-size:16px;">&#10003;</span>
<div style="font-size:12px;">
<strong style="color:#10b981;">{top_zona['zona']}:</strong> Mayor conversión ({top_zona['tasa']}%)<br>
{top_zona['cierres']} cierres de {top_zona['prospectos']} prospectos
</div></div>"""
        if not alertas_html:
            alertas_html = "<span style='font-size:13px;opacity:0.6;'>Sin alertas en este periodo.</span>"

        st.markdown(f"""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title">Alertas y Observaciones</div>
<div class="panel-subtitle">Análisis de conversión por Zona</div>
{alertas_html}
<div style="margin-top:18px;padding-top:14px;border-top:1px solid rgba(128,128,128,0.2);
            display:flex;justify-content:space-between;align-items:center;">
<span style="font-size:13px;opacity:0.65;">Conversión global</span>
<span style="font-size:20px;font-weight:800;color:#7c3aed;">{tasa_global}%</span>
</div>
</div>""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SECCIÓN: VISITAS DIARIAS / PATRÓN SEMANAL
# ═══════════════════════════════════════════════════════════════════════════════
DIAS_ES  = {0:"LUN", 1:"MAR", 2:"MIÉ", 3:"JUE", 4:"VIE", 5:"SÁB", 6:"DOM"}
DIAS_ORD = [0, 1, 2, 3, 4, 5]          # Lun-Sáb (domingos rara vez)
DIAS_COL = ["#3b82f6","#10b981","#f59e0b","#8b5cf6","#ef4444","#0ea5e9","#14b8a6"]

def render_visitas_diarias_section(df, modo):
    """Sección Visitas Diarias / Patrón Semanal — pestaña Todos."""
    df_fis = df[
        df[COL_TIPO].str.upper().isin(["PROSPECCIÓN","PROSPECCION","MANTENIMIENTO"]) &
        df.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
    ].copy()

    if df_fis.empty:
        return

    df_fis["_dow"]  = df_fis[COL_FECHA].dt.dayofweek
    df_fis["_date"] = df_fis[COL_FECHA].dt.date
    df_fis["_dom"]  = df_fis[COL_FECHA].dt.day
    total_fis = len(df_fis)

    # ── Patrón por día de semana (Lun-Sáb, sin domingos) ────────────────────
    dow_counts  = df_fis.groupby("_dow").size().reindex(range(6), fill_value=0)  # 0-5
    dias_labels = [DIAS_ES[d] for d in DIAS_ORD]
    dias_vals   = [int(dow_counts.get(d, 0)) for d in DIAS_ORD]
    # Promedio por día de semana: suma de totales Lun-Sáb / 6
    dias_sin_dom = df_fis[df_fis["_dow"] != 6]["_date"].nunique()
    avg_dia = sum(dias_vals) / 6   # referencia para el delta de cada día
    prom_diario = round(total_fis / max(dias_sin_dom, 1), 1)  # para el stat de promedio

    # ── Totales semanales (solo modo Mes) ────────────────────────────────────
    sem_totals = None
    if modo == "Mes":
        df_fis["_yw"] = (
            df_fis[COL_FECHA].dt.isocalendar()["year"].astype(str) + "-S" +
            df_fis[COL_FECHA].dt.isocalendar()["week"].astype(str).str.zfill(2)
        )
        sem_group = (df_fis.groupby("_yw").size()
                     .reset_index(name="Visitas")
                     .sort_values("_yw").reset_index(drop=True))

        # Label: rango de fechas de la semana  13/04/2026 -- 19/04/2026
        def sem_range_label(yw):
            anio, snro = yw.split("-S")
            lunes   = pd.Timestamp.fromisocalendar(int(anio), int(snro), 1)
            domingo = lunes + pd.Timedelta(days=6)
            return f"{lunes.strftime('%d/%m/%Y')} -- {domingo.strftime('%d/%m/%Y')}"

        sem_group["Label"] = sem_group["_yw"].apply(sem_range_label)
        sem_totals = sem_group

    # ── Stats compartidos (día pico, día del mes, promedio) ──────────────────
    idx_pico_dow = dias_vals.index(max(dias_vals)) if max(dias_vals) > 0 else 0
    dia_pico_nom = dias_labels[idx_pico_dow]
    dia_pico_val = dias_vals[idx_pico_dow]

    dom_counts   = df_fis[df_fis["_dow"] != 6].groupby("_dom").size()
    if not dom_counts.empty:
        dom_pico   = int(dom_counts.idxmax())
        dom_pico_v = int(dom_counts.max())
        fecha_dom  = df_fis[df_fis["_dom"] == dom_pico][COL_FECHA].iloc[0]
        dom_pico_dow = DIAS_ES[fecha_dom.dayofweek]
    else:
        dom_pico, dom_pico_v, dom_pico_dow = 0, 0, "-"



    stats_html = f"""
<div style="margin-top:12px;display:flex;flex-direction:column;gap:6px;">
<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;border-radius:6px;background:var(--secondary-background-color);">
  <div style="width:10px;height:10px;border-radius:2px;background:#3b82f6;flex-shrink:0;"></div>
  <span style="font-size:12px;flex:1;opacity:0.7;">Día pico</span>
  <span style="font-size:13px;font-weight:800;color:#3b82f6;">{dia_pico_nom} ({dia_pico_val})</span>
</div>
<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;border-radius:6px;background:var(--secondary-background-color);">
  <div style="width:10px;height:10px;border-radius:2px;background:#f59e0b;flex-shrink:0;"></div>
  <span style="font-size:12px;flex:1;opacity:0.7;">{dom_pico_dow} {dom_pico}</span>
  <span style="font-size:13px;font-weight:800;color:#f59e0b;">{dom_pico_v} visitas</span>
</div>
<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;border-radius:6px;background:var(--secondary-background-color);">
  <div style="width:10px;height:10px;border-radius:2px;background:#10b981;flex-shrink:0;"></div>
  <span style="font-size:12px;flex:1;opacity:0.7;">Promedio diario (sin domingos)</span>
  <span style="font-size:13px;font-weight:800;color:#10b981;">{prom_diario} visitas</span>
</div>
</div>"""

    # ── Cabecera ─────────────────────────────────────────────────────────────
    seccion_header(
        "📅", "Visitas Diarias / Patrón Semanal",
        "Análisis de actividad por día de semana",
        f"Total: {total_fis:,} visitas físicas",
        gradient="linear-gradient(135deg,#1e40af 0%,#2563eb 100%)"
    )

    col_izq, col_der = st.columns([1, 1]) if modo == "Mes" else [st.container(), None]

    # ── Panel izquierdo: Patrón por día ──────────────────────────────────────
    with col_izq:
        st.markdown("""
<div class="dashboard-panel" style="margin-top:12px;padding-bottom:4px;">
<div class="panel-title">Patrón de Visitas por Día</div>
<div class="panel-subtitle">Distribución diaria de actividad comercial</div>
</div>""", unsafe_allow_html=True)

        fig_dow = go.Figure()
        fig_dow.add_trace(go.Scatter(
            x=dias_labels, y=dias_vals,
            mode="lines+markers",
            line=dict(color="#3b82f6", width=3, shape="spline"),
            marker=dict(size=8, color="#3b82f6"),
            fill="tozeroy",
            fillcolor="rgba(59,130,246,0.12)",
            hovertemplate="%{x}: %{y} visitas<extra></extra>",
        ))
        fig_dow.update_layout(
            paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
            margin=dict(t=10, b=10, l=10, r=10), height=230,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=True, gridcolor="rgba(128,128,128,0.15)"),
        )
        st.plotly_chart(fig_dow, use_container_width=True, config={"displayModeBar": False})

        # Tarjetas estilo ruta-card
        cards_dow = ""
        for i, (dia, val) in enumerate(zip(dias_labels, dias_vals)):
            delta     = val - avg_dia
            delta_pct = round(delta / avg_dia * 100) if avg_dia > 0 else 0
            arrow     = "↑" if delta >= 0 else "↓"
            col_delta = "#10b981" if delta >= 0 else "#ef4444"
            color     = DIAS_COL[i]
            bar_w     = round(val / max(dias_vals) * 100) if max(dias_vals) > 0 else 0
            cards_dow += f"""
<div class="ruta-card">
<div class="ruta-header">
  <span class="ruta-name">{dia}</span>
  <span class="ruta-value" style="color:{color};">{val}</span>
</div>
<div class="ruta-progress-container">
  <div class="ruta-progress-bar">
    <div class="ruta-progress-fill" style="width:{bar_w}%;background:{color};"></div>
  </div>
  <span class="ruta-pct" style="color:{col_delta};font-size:11px;">{arrow}{'+' if delta_pct>=0 else ''}{delta_pct}%</span>
</div>
</div>"""
        st.markdown(f'<div class="rutas-container">{cards_dow}</div>', unsafe_allow_html=True)

        # Cuando el filtro es por Semana, los stats van aquí abajo
        if modo == "Semana":
            st.markdown(stats_html, unsafe_allow_html=True)

    # ── Panel derecho: Totales semanales (solo Mes) ──────────────────────────
    if modo == "Mes" and sem_totals is not None and not sem_totals.empty:
        with col_der:
            st.markdown("""
<div class="dashboard-panel" style="margin-top:12px;padding-bottom:4px;">
<div class="panel-title">Totales Semanales</div>
<div class="panel-subtitle">Visitas por semana</div>
</div>""", unsafe_allow_html=True)

            colores_W = [PALETA_RUTAS[i % len(PALETA_RUTAS)] for i in range(len(sem_totals))]
            fig_sem = go.Figure(go.Bar(
                x=sem_totals["Label"], y=sem_totals["Visitas"],
                marker_color=colores_W,
                text=sem_totals["Visitas"], textposition="outside",
                cliponaxis=False,
            ))
            fig_sem.update_layout(
                paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                margin=dict(t=10, b=30, l=10, r=10), height=250,
                xaxis=dict(showgrid=False, tickangle=-30, tickfont=dict(size=9)),
                yaxis=dict(showgrid=True, gridcolor="rgba(128,128,128,0.15)"),
            )
            st.plotly_chart(fig_sem, use_container_width=True, config={"displayModeBar": False})

            # Stats sólo en modo Mes (van aquí bajo el gráfico semanal)
            st.markdown(stats_html, unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SECCIÓN: INDICADORES DE PROSPECCIÓN Y MANTENIMIENTO
# ═══════════════════════════════════════════════════════════════════════════════
def render_indicadores_giro_section(df, region_nombre):
    """Sección Indicadores de Prospección y Mantenimiento: distribución y efectividad por giro."""
    if COL_TIPO_CLI not in df.columns:
        return
    if df.empty:
        return

    orden_etapa = {e: i for i, e in enumerate(ETAPAS_EMBUDO)}

    # ── Cabecera ───────────────────────────────────────────────────────────
    st.markdown('<div style="margin-top:30px;"></div>', unsafe_allow_html=True)
    seccion_header(
        "📈", "Indicadores de Prospección y Mantenimiento",
        f"Análisis por tipo de cliente (giro) · {region_nombre}",
        f"{df[COL_TIPO_CLI].nunique()} giros registrados",
        gradient="linear-gradient(135deg,#0f172a 0%,#1e40af 100%)"
    )

    # ── Datos por giro ─────────────────────────────────────────────────────────
    giros = sorted(df[COL_TIPO_CLI].dropna().unique().tolist())
    # Prospección: solo para cálculo conversion
    df_pros_gbl = df[df[COL_TIPO].str.upper().isin(["PROSPECCIÓN", "PROSPECCION"])]
    # Motivo TOMAR PEDIDO
    MOTIVO_TP = "TOMAR PEDIDO"

    giro_data = []  # dict por giro con todos los stats
    for i, g in enumerate(giros):
        df_g = df[df[COL_TIPO_CLI] == g]
        vis_total  = len(df_g)
        clientes_u = df_g[COL_CLIENTE].nunique()
        # Prospectos únicos (solo Prospección)
        df_g_pros = df_pros_gbl[df_pros_gbl[COL_TIPO_CLI] == g]
        prospectos_u = df_g_pros[COL_CLIENTE].nunique()
        # Cierres
        n_cierres = 0
        if COL_MOTIVO in df_g_pros.columns and not df_g_pros.empty:
            valid = df_g_pros[df_g_pros[COL_MOTIVO].str.upper().isin([e.upper() for e in ETAPAS_EMBUDO])].copy()
            if not valid.empty:
                valid["_ord"] = valid[COL_MOTIVO].str.upper().map(orden_etapa)
                ultima = valid.sort_values("_ord").groupby(COL_CLIENTE).last()[[COL_MOTIVO]].reset_index()
                n_cierres = ultima[ultima[COL_MOTIVO].str.upper() == "CIERRE"].shape[0]
        tasa_conv = round(n_cierres / prospectos_u * 100, 1) if prospectos_u > 0 else 0.0
        # Visitas de PROSPECCIÓN para la dona
        visitas_pros = len(df_g_pros)
        # Visitas de MANTENIMIENTO para efectividad
        df_g_mant = df[df[COL_TIPO_CLI] == g]
        df_g_mant = df_g_mant[df_g_mant[COL_TIPO].str.upper() == "MANTENIMIENTO"]
        vis_mant     = len(df_g_mant)
        clientes_mant = df_g_mant[COL_CLIENTE].nunique()
        if COL_MOTIVO in df_g_mant.columns and vis_mant > 0:
            n_tp = df_g_mant[df_g_mant[COL_MOTIVO].str.upper() == MOTIVO_TP].shape[0]
        else:
            n_tp = 0
        efect = round(n_tp / vis_mant * 100, 1) if vis_mant > 0 else 0.0
        giro_data.append({
            "giro": g, "visitas": vis_total, "visitas_pros": visitas_pros,
            "visitas_mant": vis_mant, "clientes_mant": clientes_mant,
            "clientes_u": clientes_u,
            "prospectos_u": prospectos_u, "cierres": n_cierres,
            "tasa_conv": tasa_conv, "tomar_pedido": n_tp, "efect": efect,
            "color": PALETA_RUTAS[i % len(PALETA_RUTAS)]
        })

    total_vis = sum(g["visitas"] for g in giro_data)
    total_pros_dona = sum(g["visitas_pros"] for g in giro_data)
    if not giro_data or total_vis == 0:
        st.info("Sin datos para mostrar indicadores por giro.")
        return

    col_dona, col_efect = st.columns([1, 1])

    # ── Indicador 1: Distribución por tipo de cliente (dona) ────────────────
    with col_dona:
        st.markdown("""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title">Distribución por Tipo de Cliente</div>
<div class="panel-subtitle">Visitas de Prospección por giro de cliente</div>
</div>""", unsafe_allow_html=True)

        fig_dona = go.Figure(go.Pie(
            labels=[g["giro"] for g in giro_data],
            values=[g["visitas_pros"] for g in giro_data],
            hole=0.52,
            marker_colors=[g["color"] for g in giro_data],
            textinfo="label+percent",
            hovertemplate="%{label}: %{value} visitas (%{percent})<extra></extra>",
        ))
        fig_dona.update_layout(
            paper_bgcolor="rgba(0,0,0,0)",
            margin=dict(t=10, b=130, l=10, r=10),
            height=420,
            legend=dict(orientation="h", yanchor="bottom", y=-0.38, xanchor="center", x=0.5),
            showlegend=True,
        )
        st.plotly_chart(fig_dona, use_container_width=True, config={"displayModeBar": False})

        # Tarjetas por giro: visitas prospección | prospectos únicos | cierres | tasa
        cards_giro = ""
        for g in sorted(giro_data, key=lambda x: x["visitas_pros"], reverse=True):
            color = g["color"]
            bar_w = round(g["visitas_pros"] / total_pros_dona * 100) if total_pros_dona > 0 else 0
            cards_giro += f"""
<div class="ruta-card">
<div class="ruta-header">
  <span class="ruta-name">{g['giro']}</span>
  <span class="ruta-value" style="color:{color};">{g['visitas_pros']}</span>
</div>
<div class="ruta-progress-container">
  <div class="ruta-progress-bar">
    <div class="ruta-progress-fill" style="width:{bar_w}%;background:{color};"></div>
  </div>
</div>
<div style="display:flex;gap:16px;margin-top:6px;font-size:11px;opacity:0.75;">
  <span>&#128100; {g['prospectos_u']} prospec.</span>
  <span>&#9989; {g['cierres']} cierres</span>
  <span>&#128200; {g['tasa_conv']}% conv.</span>
</div>
</div>"""
        st.markdown(f'<div class="rutas-container" style="margin-top:12px;">{cards_giro}</div>',
                    unsafe_allow_html=True)

    # ── Indicador 2: Efectividad por giro de cartera ──────────────────────
    with col_efect:
        st.markdown("""
<div class="dashboard-panel" style="margin-bottom:15px;">
<div class="panel-title">Efectividad por Giro de Cartera</div>
<div class="panel-subtitle">Tomar Pedido sobre visitas de Mantenimiento por giro</div>
</div>""", unsafe_allow_html=True)

        giro_efect_sorted = sorted(giro_data, key=lambda x: x["efect"], reverse=True)
        for g in giro_efect_sorted:
            color  = g["color"]
            efect  = g["efect"]
            bar_w  = min(efect, 100)
            st.markdown(f"""
<div class="ruta-card" style="margin-bottom:10px;">
<div class="ruta-header">
  <span class="ruta-name" style="font-size:14px;font-weight:700;">{g['giro']}</span>
  <span class="ruta-value" style="color:{color};font-size:16px;">{efect}%</span>
</div>
<div class="ruta-progress-container">
  <div class="ruta-progress-bar">
    <div class="ruta-progress-fill" style="width:{bar_w}%;background:{color};"></div>
  </div>
</div>
<div style="display:flex;gap:18px;margin-top:8px;font-size:11px;opacity:0.75;">
  <span>&#128200; {g['visitas_mant']} vis. mant.</span>
  <span>&#128100; {g['clientes_mant']} clientes</span>
  <span>&#128722; {g['tomar_pedido']} Tomar Pedido</span>
</div>
</div>""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# PESTAÑAS
# ═══════════════════════════════════════════════════════════════════════════════
tab_todos, tab_prov, tab_lima = st.tabs(["🌎 Todos", "🏞️ Provincia", "🏙️ Lima"])

with tab_todos:
    if dff.empty:
        st.info("No hay datos para mostrar.")
    else:
        render_region_dashboard(dff, "Todos", is_todos=True)
        render_visitas_diarias_section(dff, modo_fecha)
        render_indicadores_giro_section(dff, "Todos")

with tab_prov:
    if COL_REGION in dff.columns:
        df_prov = dff[dff[COL_REGION].str.upper() == "PROVINCIA"]
        if df_prov.empty:
            st.info("Sin datos de Provincia para los filtros actuales.")
        else:
            render_region_dashboard(df_prov, "Provincia")
        render_conversion_section(df_prov, "Provincia")
        render_indicadores_giro_section(df_prov, "Provincia")
    else:
        st.warning("Falta la columna Región.")

with tab_lima:
    if COL_REGION in dff.columns:
        df_lim = dff[dff[COL_REGION].str.upper() == "LIMA"]
        if df_lim.empty:
            st.info("Sin datos de Lima para los filtros actuales.")
        else:
            render_region_dashboard(df_lim, "Lima")
        render_conversion_section(df_lim, "Lima")
        render_indicadores_giro_section(df_lim, "Lima")
    else:
        st.warning("Falta la columna Región.")

# ═══════════════════════════════════════════════════════════════════════════════
# SIDEBAR – EXPORTAR
# ═══════════════════════════════════════════════════════════════════════════════
with st.sidebar:
    st.divider()
    st.markdown("### 📥 Exportar")
    if st.button("Presentación PPTX", use_container_width=True):

        try:
            import io
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            import matplotlib.patches as mpatches

            # ── Paleta y helpers ──────────────────────────────────────────────
            AZUL      = RGBColor(0x1d, 0x4e, 0xd8)
            VERDE     = RGBColor(0x06, 0x5f, 0x46)
            BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
            GRIS_OSC  = RGBColor(0x1e, 0x3a, 0x8a)
            GRIS_CLR  = RGBColor(0x64, 0x74, 0x8B)
            W, H = Inches(13.33), Inches(7.5)

            def set_bg(slide, r, g, b):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)

            def txbox(slide, text, left, top, width, height,
                      font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.word_wrap = wrap
                p  = tf.paragraphs[0]
                p.alignment = align
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = color
                return tb

            def add_rect(slide, left, top, width, height, fill_rgb, radius=False):
                shape = slide.shapes.add_shape(
                    1, left, top, width, height)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
                shape.line.fill.background()
                return shape

            def fig_to_bytes(fig):
                buf = io.BytesIO()
                fig.savefig(buf, format="png", bbox_inches="tight", dpi=150, facecolor=fig.get_facecolor())
                buf.seek(0)
                return buf

            def bar_chart_bytes(grupos_df, title, palette):
                """Dibuja un gráfico de barras con matplotlib y devuelve bytes PNG."""
                fig, ax = plt.subplots(figsize=(9, 3.5), facecolor="#F8FAFC")
                ax.set_facecolor("#F8FAFC")
                zonas  = grupos_df[COL_ZONA].tolist()
                visits = grupos_df["Visitas"].tolist()
                colors = [palette[i % len(palette)] for i in range(len(zonas))]
                bars = ax.bar(zonas, visits, color=colors, edgecolor="none", width=0.55)
                for bar, v in zip(bars, visits):
                    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                            str(v), ha="center", va="bottom", fontsize=9, fontweight="bold", color="#1e3a8a")
                ax.set_title(title, fontsize=11, fontweight="bold", color="#1e3a8a", pad=8)
                ax.set_xlabel("")
                ax.spines[["top","right","left"]].set_visible(False)
                ax.tick_params(axis="y", left=False, labelleft=False)
                ax.tick_params(axis="x", labelsize=8)
                plt.tight_layout()
                return fig_to_bytes(fig)

            def kpi_slide_data(df_r, region_nombre, prs):
                """Agrega una slide con KPIs + tabla de zonas para la región."""
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                set_bg(slide, 0xF8, 0xFA, 0xFC)

                # Cabecera azul
                add_rect(slide, 0, 0, W, Inches(1.1), (0x1d, 0x4e, 0xd8))
                txbox(slide, f"📊  Resumen Ejecutivo · {region_nombre}",
                      Inches(0.3), Inches(0.15), Inches(9), Inches(0.5),
                      font_size=22, bold=True, color=BLANCO)
                txbox(slide, rango_label,
                      Inches(0.3), Inches(0.65), Inches(9), Inches(0.35),
                      font_size=12, bold=False, color=BLANCO)

                # KPIs
                tot_vis, tot_pros, n_cierres, t_conv = calc_kpis(df_r)
                kpi_data = [
                    ("Visitas Totales", str(tot_vis),  (0x3b, 0x82, 0xf6)),
                    ("Prospectos",      str(tot_pros), (0x10, 0xb9, 0x81)),
                    ("Cierres",         str(n_cierres),(0xf5, 0x9e, 0x0b)),
                    ("Conversión",      f"{t_conv}%",  (0x8b, 0x5c, 0xf6)),
                ]
                card_w = Inches(2.8)
                card_h = Inches(1.4)
                card_top = Inches(1.25)
                gap = Inches(0.4)
                for i, (label, val, rgb) in enumerate(kpi_data):
                    left = Inches(0.4) + i * (card_w + gap)
                    add_rect(slide, left, card_top, card_w, card_h, (0xFF,0xFF,0xFF))
                    # Borde izquierdo coloreado simulado con rectángulo estrecho
                    add_rect(slide, left, card_top, Inches(0.07), card_h, rgb)
                    txbox(slide, label, left + Inches(0.15), card_top + Inches(0.12),
                          card_w - Inches(0.2), Inches(0.35), font_size=10, color=GRIS_CLR)
                    txbox(slide, val,   left + Inches(0.15), card_top + Inches(0.5),
                          card_w - Inches(0.2), Inches(0.6), font_size=26, bold=True,
                          color=RGBColor(*rgb))

                # Tabla de zonas (si hay datos)
                if COL_ZONA in df_r.columns:
                    df_fis = df_r[
                        df_r[COL_TIPO].str.upper().isin(["PROSPECCIÓN","PROSPECCION","MANTENIMIENTO"]) &
                        df_r.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
                    ]
                    grupos = df_fis.groupby(COL_ZONA).size().reset_index(name="Visitas").sort_values("Visitas", ascending=False).reset_index(drop=True)
                    total  = grupos["Visitas"].sum()

                    tbl_top  = Inches(2.85)
                    tbl_left = Inches(0.4)
                    tbl_w    = Inches(12.5)
                    row_h    = Inches(0.32)
                    cols_w   = [Inches(3.5), Inches(2), Inches(2), Inches(4.5)]
                    headers  = ["Zona", "Visitas", "% Part.", "Estado"]

                    # Cabecera tabla
                    x = tbl_left
                    for ci, (hdr, cw) in enumerate(zip(headers, cols_w)):
                        add_rect(slide, x, tbl_top, cw, row_h, (0x1d, 0x4e, 0xd8))
                        txbox(slide, hdr, x + Inches(0.05), tbl_top + Inches(0.04),
                              cw - Inches(0.1), row_h, font_size=10, bold=True, color=BLANCO)
                        x += cw

                    # Filas
                    EST_COLORS = {"ACTIVO": (0x10,0xb9,0x81), "REGULAR": (0xf5,0x9e,0x0b), "BAJO": (0xef,0x44,0x44)}
                    for ri, row in grupos.iterrows():
                        y = tbl_top + (ri+1)*row_h
                        if y + row_h > Inches(7.1):
                            break
                        zona  = row[COL_ZONA]
                        vis   = row["Visitas"]
                        pct   = round(vis/total*100, 1) if total > 0 else 0
                        est   = obtener_estado(vis, df_estado_usar, num_periodos)
                        est_u = est.upper()
                        c_est = EST_COLORS.get(next((k for k in EST_COLORS if k in est_u), ""), (0x64,0x74,0x8B))
                        bg_row = (0xFF,0xFF,0xFF) if ri % 2 == 0 else (0xF1,0xF5,0xF9)
                        x = tbl_left
                        for ci, (val, cw) in enumerate(zip([zona, str(vis), f"{pct}%", est], cols_w)):
                            add_rect(slide, x, y, cw, row_h, bg_row)
                            col_txt = RGBColor(*c_est) if ci == 3 else GRIS_OSC
                            txbox(slide, val, x + Inches(0.05), y + Inches(0.04),
                                  cw - Inches(0.1), row_h, font_size=9,
                                  bold=(ci == 3), color=col_txt)
                            x += cw

            def mant_slide(df_r, region_nombre, prs):
                """Agrega slide con gráfico de barras de Mantenimiento."""
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                set_bg(slide, 0xF8, 0xFA, 0xFC)
                add_rect(slide, 0, 0, W, Inches(1.1), (0x06, 0x5f, 0x46))
                txbox(slide, f"🔧  Visitas por Zona - Mantenimiento · {region_nombre}",
                      Inches(0.3), Inches(0.15), Inches(10), Inches(0.5),
                      font_size=20, bold=True, color=BLANCO)
                txbox(slide, rango_label, Inches(0.3), Inches(0.65), Inches(9), Inches(0.35),
                      font_size=12, color=BLANCO)

                df_mant_fis = df_r[
                    (df_r[COL_TIPO].str.upper() == "MANTENIMIENTO") &
                    df_r.get(COL_TIPO_VIS, pd.Series(dtype=str)).str.upper().isin(["FÍSICA","FISICA"])
                ]
                grupos_m = df_mant_fis.groupby(COL_ZONA).size().reset_index(name="Visitas").sort_values("Visitas", ascending=False).reset_index(drop=True)
                if grupos_m.empty:
                    txbox(slide, "Sin datos de Mantenimiento físico para el periodo.",
                          Inches(1), Inches(2), Inches(10), Inches(1), font_size=14, color=GRIS_CLR)
                    return

                chart_bytes = bar_chart_bytes(grupos_m, "Distribución de Visitas por Zona (Mantenimiento)", PALETA_RUTAS)
                slide.shapes.add_picture(chart_bytes, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5))

            # ── Armar presentación ─────────────────────────────────────────────
            prs = Presentation()
            prs.slide_width  = W
            prs.slide_height = H

            # Slide portada
            slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide_cover, 0x1d, 0x4e, 0xd8)
            txbox(slide_cover, "🗺️  Reporte de Visitas Comerciales",
                  Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                  font_size=36, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
            txbox(slide_cover, rango_label.capitalize(),
                  Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                  font_size=18, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)
            txbox(slide_cover, f"Vendedor: {sel_vendedor}",
                  Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                  font_size=13, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

            # Slides por región
            regiones = [("Todos", dff)]
            if COL_REGION in dff.columns:
                df_prov = dff[dff[COL_REGION].str.upper() == "PROVINCIA"]
                df_lim  = dff[dff[COL_REGION].str.upper() == "LIMA"]
                if not df_prov.empty:
                    regiones.append(("Provincia", df_prov))
                if not df_lim.empty:
                    regiones.append(("Lima", df_lim))

            for nombre, df_r in regiones:
                kpi_slide_data(df_r, nombre, prs)
                if nombre != "Todos":
                    mant_slide(df_r, nombre, prs)

            # Guardar en buffer y ofrecer descarga
            buf_pptx = io.BytesIO()
            prs.save(buf_pptx)
            buf_pptx.seek(0)
            st.sidebar.download_button(
                label="⬇️ Descargar PPTX",
                data=buf_pptx,
                file_name=f"reporte_visitas_{rango_label.replace(' ','_')[:40]}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            st.sidebar.success("✅ Presentación lista para descargar.")

        except Exception as e:
            st.sidebar.error(f"Error al generar PPTX: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# DETALLE DE VISITAS
# ═══════════════════════════════════════════════════════════════════════════════
st.markdown("## 📋 Detalle de Visitas")
cols_mostrar = [COL_FECHA, COL_VENDEDOR, COL_ZONA, COL_REGION, COL_TIPO, COL_TIPO_VIS, COL_CLIENTE, COL_MOTIVO, COL_RESULTADO]

def prep_detalle(df):
    disponibles = [c for c in cols_mostrar if c in df.columns]
    d = df[disponibles].copy()
    if COL_FECHA in d.columns:
        d[COL_FECHA] = d[COL_FECHA].dt.date
        d = d.sort_values(COL_FECHA, ascending=False)
    return d

tdt, tdp, tdl = st.tabs(["🌎 Todos", "🏞️ Provincia", "🏙️ Lima"])
with tdt:
    st.dataframe(prep_detalle(dff), use_container_width=True, hide_index=True)
with tdp:
    if COL_REGION in dff.columns:
        st.dataframe(prep_detalle(dff[dff[COL_REGION].str.upper() == "PROVINCIA"]), use_container_width=True, hide_index=True)
with tdl:
    if COL_REGION in dff.columns:
        st.dataframe(prep_detalle(dff[dff[COL_REGION].str.upper() == "LIMA"]), use_container_width=True, hide_index=True)
